
import io
import os
import re
import json
import base64
import subprocess
import tempfile
from pathlib import Path

import requests
import streamlit as st
from pptx import Presentation
from docx import Document
from docx.shared import Pt
from docx.oxml.ns import qn
from openai import OpenAI

st.set_page_config(
    page_title="PPT教学设计生成器",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded",
)

# -----------------------------
# 基础配置
# -----------------------------
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "").strip()
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY", "").strip()
APP_PASSWORD = os.getenv("APP_PASSWORD", "").strip()
DEFAULT_VISION_MODEL = os.getenv("VISION_MODEL", "gpt-5.6").strip()
DEFAULT_DEEPSEEK_MODEL = os.getenv("DEEPSEEK_MODEL", "deepseek-v4-flash").strip()
DEFAULT_OPENAI_WRITER_MODEL = os.getenv("OPENAI_WRITER_MODEL", "gpt-5.6").strip()

MAX_PPT_MB = int(os.getenv("MAX_PPT_MB", "45"))
MAX_DOCX_MB = int(os.getenv("MAX_DOCX_MB", "15"))

SECTION_KEYWORDS = [
    "课题", "课型", "教材分析", "学情分析", "教学目标",
    "教学重难点", "教学重点", "教学难点", "教学方法", "教学资源",
    "教学过程", "课堂小结", "课后作业", "作业设计",
    "板书设计", "教学反思", "预设教学反思"
]

VISION_PROMPT = """
你是一名初中社会学科PPT课堂设计分析助手。
请逐页分析这份PPT，必须同时利用页面文字与页面视觉信息。

重点识别：
1. 图片、人物、历史场景、文物图片及其在教学中的作用；
2. 历史地图/地理地图中的位置、路线、范围、方向、空间变化；
3. 图表、时间轴、流程图、对比表、示意图中的关系；
4. 史料截图中的可见信息；
5. PPT明确出现的任务、问题、探究、合作学习、练习、小结；
6. 页面之间的教学顺序和任务结构。

严格规则：
- 只依据PPT可见内容，不补充课件没有呈现的材料。
- 不要臆测图片中无法确认的细节。
- 如果页面主要是装饰图，请明确说明。
- 课堂活动线索必须能在PPT中找到依据。

请按页码输出“PPT视觉识读报告”，每页包含：
【页码】
【页面主题】
【文字要点】
【图片/地图/图表/史料信息】
【教学功能】
【可进入教案的教师活动】
【可进入教案的学生活动】
"""

WRITER_SYSTEM = """
你是一名初中社会学科教学设计编辑器。
你的任务不是重新创作一份脱离底稿的教案，而是：
“以原教学设计为底稿，在原框架上校正，使其与PPT完全适配并可直接使用”。

最高优先级规则：
1. 原教学设计的栏目名称、栏目顺序和基本框架必须保留。
2. PPT是课堂内容、任务结构和教学顺序的最高依据。
3. 原教案中凡与PPT不匹配、PPT没有出现、改变PPT任务顺序、虚构材料或活动的内容，必须删除。
4. 原教案中与PPT一致且表述合理的内容应保留，避免无意义重写。
5. PPT中真实存在但原教案遗漏的任务、问题、史料、地图、图片、图表、时间轴、课堂检测等，应补入对应环节。
6. 必须结合PPT视觉识读报告，不能只依据文本框。
7. 不得虚构PPT未展示的史料、数据、图片信息、教材页码、学生回答或课堂效果。
8. 教师活动、学生活动、设计意图必须形成明确对应，并能从PPT找到依据。
9. PPT已有“任务一/任务二/任务三”等结构时，教学过程必须保留相同结构与顺序。
10. 如果没有真实课堂反馈，教学反思只能写成预设性的反思内容，不能声称课堂已经发生。
11. 输出必须是教师可直接使用的完成稿，不写“建议教师补充……”等说明性废话；确实无依据的项目可写“需教师补充”。
12. 每个原栏目都必须返回最终内容，即使无需修改。
13. 只返回合法 JSON，不要代码围栏。

JSON：
{
  "title": "课题名称",
  "sections": [
    {"heading": "原模板栏目名称", "content": "最终可直接使用的正文"}
  ],
  "alignment_check": {
    "removed": ["被删除的原教案不匹配内容"],
    "kept": ["保留的匹配内容"],
    "added": ["根据PPT补入的内容"],
    "warnings": ["仍需人工确认的事项"]
  }
}
"""

# -----------------------------
# 登录
# -----------------------------
def require_login():
    if not APP_PASSWORD:
        return
    if st.session_state.get("authenticated"):
        return
    st.title("📚 PPT教学设计生成器")
    st.caption("受保护的教师备课工具")
    pwd = st.text_input("访问密码", type="password")
    if st.button("登录", type="primary"):
        if pwd == APP_PASSWORD:
            st.session_state["authenticated"] = True
            st.rerun()
        else:
            st.error("密码错误")
    st.stop()

require_login()

# -----------------------------
# 文件解析
# -----------------------------
def validate_size(upload, max_mb, label):
    if upload is None:
        return
    size = len(upload.getvalue())
    if size > max_mb * 1024 * 1024:
        raise ValueError(f"{label} 超过 {max_mb}MB 限制。")

def extract_ppt_text(ppt_bytes):
    prs = Presentation(io.BytesIO(ppt_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and getattr(shape, "text", "").strip():
                texts.append(shape.text.strip())
        slides.append({"page": i, "text": "\n".join(texts)})
    return slides

def extract_template(docx_bytes):
    doc = Document(io.BytesIO(docx_bytes))
    raw = []
    headings = []

    def maybe_heading(t):
        clean = re.sub(r"\s+", "", t)
        for k in SECTION_KEYWORDS:
            if clean == k or clean.startswith(k + "（") or clean.startswith(k + "("):
                if k not in headings:
                    headings.append(k)
                return

    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            raw.append(t)
            maybe_heading(t)

    for ti, table in enumerate(doc.tables, 1):
        raw.append(f"[模板表格{ti}]")
        for row in table.rows:
            vals = [c.text.strip() for c in row.cells]
            raw.append(" | ".join(vals))
            for v in vals:
                if v:
                    maybe_heading(v)

    whole = "\n".join(raw)
    if len(headings) < 5:
        headings = [k for k in SECTION_KEYWORDS if k in whole]

    return whole, headings

# -----------------------------
# 云端PPT→PDF
# Docker内安装LibreOffice
# -----------------------------
def pptx_to_pdf(ppt_bytes):
    with tempfile.TemporaryDirectory(prefix="pptdesign_") as td:
        td = Path(td)
        ppt = td / "slides.pptx"
        ppt.write_bytes(ppt_bytes)
        cmd = [
            "libreoffice", "--headless", "--nologo", "--nodefault",
            "--nolockcheck", "--nofirststartwizard",
            "--convert-to", "pdf",
            "--outdir", str(td),
            str(ppt)
        ]
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
        pdf = td / "slides.pdf"
        if proc.returncode != 0 or not pdf.exists():
            raise RuntimeError(
                "PPT转PDF失败。服务器LibreOffice返回："
                + (proc.stderr[-1000:] or proc.stdout[-1000:])
            )
        return pdf.read_bytes()

# -----------------------------
# OpenAI视觉识读
# PDF会同时解析文字与页图
# -----------------------------
def analyze_pdf_visual(pdf_bytes, model, api_key):
    if not api_key:
        raise RuntimeError("请先在网页 API 设置中填写 OpenAI API Key。")
    if len(pdf_bytes) > 50 * 1024 * 1024:
        raise ValueError("转换后的PDF超过50MB，无法一次提交视觉分析。")

    client = OpenAI(api_key=api_key)
    b64 = base64.b64encode(pdf_bytes).decode("utf-8")
    response = client.responses.create(
        model=model,
        input=[{
            "role": "user",
            "content": [
                {
                    "type": "input_file",
                    "filename": "slides.pdf",
                    "file_data": f"data:application/pdf;base64,{b64}",
                    "detail": "high",
                },
                {"type": "input_text", "text": VISION_PROMPT},
            ],
        }],
    )
    return response.output_text

# -----------------------------
# 写作模型
# -----------------------------
def get_deepseek_models(api_key):
    if not api_key:
        return []
    try:
        r = requests.get(
            "https://api.deepseek.com/models",
            headers={"Authorization": f"Bearer {api_key}"},
            timeout=10,
        )
        r.raise_for_status()
        return [x.get("id") for x in r.json().get("data", []) if x.get("id")]
    except Exception:
        return []

def build_writer_prompt(template_text, headings, ppt_text, vision_report,
                        subject, grade, textbook, periods, extra):
    framework = "\n".join(f"{i+1}. {h}" for i, h in enumerate(headings))
    return f"""
【基本信息】
学科：{subject}
年级：{grade}
教材版本：{textbook or "未提供"}
课时：{periods}
特殊要求：{extra or "无"}

【原教学设计完整文本】
{template_text}

【必须保留的原教学设计框架及顺序】
{framework}

【PPT逐页文本】
{ppt_text}

【PPT整页视觉识读报告】
{vision_report}

请执行“原教案框架校正”：

A. 逐项核对原教案：内容是否能在PPT文字或视觉报告中找到依据。
B. 删除：所有与PPT不匹配的旧流程、旧史料、旧活动、旧问题。
C. 保留：与PPT一致且已经写好的有效内容。
D. 补充：PPT真实存在但原教案遗漏的任务、图片、地图、史料、图表、检测。
E. 重写教学过程时严格遵循PPT真实任务顺序。
F. 教学过程必须具体写清教师活动、学生活动、设计意图。
G. 地图、图片、图表、史料、时间轴必须写出“怎么用于课堂”，但不得超出视觉报告可确认信息。
H. 最终文本必须直接可用于教师备课，不输出分析过程。

请严格按系统要求返回JSON。
"""

def call_deepseek(model, prompt, api_key, thinking=True):
    if not api_key:
        raise RuntimeError("请先在网页 API 设置中填写 DeepSeek API Key。")
    client = OpenAI(api_key=api_key, base_url="https://api.deepseek.com")
    extra_body = {"thinking": {"type": "enabled" if thinking else "disabled"}}
    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": WRITER_SYSTEM},
            {"role": "user", "content": prompt + "\n请返回合法JSON对象。"},
        ],
        response_format={"type": "json_object"},
        reasoning_effort="high",
        extra_body=extra_body,
    )
    return response.choices[0].message.content

def call_openai_writer(model, prompt, api_key):
    if not api_key:
        raise RuntimeError("请先在网页 API 设置中填写 OpenAI API Key。")
    client = OpenAI(api_key=api_key)
    response = client.responses.create(
        model=model,
        instructions=WRITER_SYSTEM,
        input=prompt + "\n只返回合法JSON对象。",
    )
    return response.output_text

def safe_json_load(text):
    t = text.strip()
    if t.startswith("```"):
        t = re.sub(r"^```(?:json)?\s*", "", t)
        t = re.sub(r"\s*```$", "", t)
    return json.loads(t)

# -----------------------------
# Word输出
# 尽量继承原文档页面设置/页眉页脚/样式
# -----------------------------
def clear_body_keep_sections(doc):
    body = doc._element.body
    sect_pr = body.sectPr
    for child in list(body):
        if child is not sect_pr:
            body.remove(child)

def east_asia_font(run, font="宋体"):
    run.font.name = font
    run._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), font)

def make_final_docx(original_docx_bytes, result, headings, subject, grade, periods):
    # 打开原文档，保留其主题、样式、页边距、页眉页脚，然后重建正文。
    doc = Document(io.BytesIO(original_docx_bytes))
    clear_body_keep_sections(doc)

    title_text = result.get("title") or "教学设计"
    p = doc.add_paragraph()
    p.alignment = 1
    r = p.add_run(title_text)
    r.bold = True
    r.font.size = Pt(16)
    east_asia_font(r, "宋体")

    section_map = {}
    for item in result.get("sections", []):
        heading = str(item.get("heading", "")).strip()
        if heading:
            section_map[heading] = str(item.get("content", "")).strip()

    # 原框架顺序绝不改变
    for h in headings:
        hp = doc.add_heading(h, level=1)
        for rr in hp.runs:
            east_asia_font(rr, "黑体")
        content = section_map.get(h, "").strip() or "需教师补充"
        for block in content.splitlines():
            s = block.strip()
            if not s:
                continue
            p = doc.add_paragraph()
            r = p.add_run(s)
            east_asia_font(r, "宋体")
            r.font.size = Pt(10.5)

    p = doc.add_paragraph()
    r = p.add_run(f"学科：{subject}    年级：{grade}    课时：{periods}")
    east_asia_font(r, "宋体")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()

def check_markdown(result):
    c = result.get("alignment_check", {})
    groups = [
        ("已删除的不匹配内容", "removed"),
        ("保留的有效内容", "kept"),
        ("根据PPT补入的内容", "added"),
        ("仍需人工确认", "warnings"),
    ]
    out = []
    for title, key in groups:
        out.append(f"### {title}")
        vals = c.get(key) or []
        if vals:
            out.extend([f"- {x}" for x in vals])
        else:
            out.append("- 无")
    return "\n".join(out)

# -----------------------------
# UI
# -----------------------------
st.title("📚 PPT教学设计一键生成器")
st.caption("云端视觉版 · 原教案框架内校正 · 删除不匹配内容 · 可直接下载Word")

with st.sidebar:
    st.header("🔑 API 设置")
    st.caption("可直接在网页填写。Key仅保存在当前网页会话内；刷新/关闭页面后需要重新填写。服务器环境变量仍可作为备用。")

    if "openai_key" not in st.session_state:
        st.session_state.openai_key = OPENAI_API_KEY
    if "deepseek_key" not in st.session_state:
        st.session_state.deepseek_key = DEEPSEEK_API_KEY

    openai_key = st.text_input(
        "OpenAI API Key（视觉识读）",
        value=st.session_state.openai_key,
        type="password",
        placeholder="sk-..."
    )
    deepseek_key = st.text_input(
        "DeepSeek API Key（教案生成）",
        value=st.session_state.deepseek_key,
        type="password",
        placeholder="sk-..."
    )
    st.session_state.openai_key = openai_key.strip()
    st.session_state.deepseek_key = deepseek_key.strip()

    c1, c2 = st.columns(2)
    with c1:
        if st.button("测试 OpenAI"):
            try:
                OpenAI(api_key=st.session_state.openai_key).models.list()
                st.success("OpenAI 连接成功")
            except Exception as e:
                st.error(f"连接失败：{e}")
    with c2:
        if st.button("测试 DeepSeek"):
            models = get_deepseek_models(st.session_state.deepseek_key)
            if models:
                st.success("DeepSeek 连接成功")
            else:
                st.error("连接失败，请检查 Key 或网络")

    st.write("视觉识读：", "✅ 已配置" if st.session_state.openai_key else "❌ 未配置")
    st.write("DeepSeek：", "✅ 已配置" if st.session_state.deepseek_key else "❌ 未配置")

    st.divider()
    st.header("教学设置")
    subject = st.selectbox("学科", ["历史", "地理", "道德与法治", "社会"])
    grade = st.text_input("年级", "八年级")
    textbook = st.text_input("教材版本", "统编版")
    periods = st.text_input("课时", "1课时")
    extra = st.text_area("特殊要求", "")

    st.divider()
    st.header("模型设置")
    vision_model = st.text_input("视觉模型", DEFAULT_VISION_MODEL)

    writer_options = []
    if st.session_state.deepseek_key:
        writer_options.append("DeepSeek")
    if st.session_state.openai_key:
        writer_options.append("OpenAI")
    if not writer_options:
        writer_options = ["未配置"]

    writer_provider = st.selectbox("教案生成模型", writer_options)

    if writer_provider == "DeepSeek":
        live_models = get_deepseek_models(st.session_state.deepseek_key)
        model_options = live_models or [DEFAULT_DEEPSEEK_MODEL, "deepseek-chat", "deepseek-reasoner"]
        model_options = list(dict.fromkeys(model_options))
        writer_model = st.selectbox("DeepSeek模型", model_options)
        thinking = st.toggle("Thinking模式", value=True)
    elif writer_provider == "OpenAI":
        writer_model = st.text_input("OpenAI写作模型", DEFAULT_OPENAI_WRITER_MODEL)
        thinking = False
    else:
        writer_model = ""
        thinking = False

col1, col2 = st.columns(2)
with col1:
    ppt_file = st.file_uploader(
        "① 上传本课PPT",
        type=["pptx"],
        help=f"建议不超过{MAX_PPT_MB}MB"
    )
with col2:
    docx_file = st.file_uploader(
        "② 上传原教学设计",
        type=["docx"],
        help="应上传格式正确、但教学流程需要与PPT校正的原教学设计"
    )

st.markdown("""
**本版本的处理原则：**
- 保留原教学设计栏目和顺序；
- 删除所有与PPT不匹配的旧内容；
- 识读PPT中的文字、地图、图片、图表、史料截图和时间轴；
- 只补入PPT真实支持的教学内容；
- 输出一份可直接使用的Word教学设计。
""")

if st.button("🚀 开始生成最终教学设计", type="primary", use_container_width=True):
    if not st.session_state.get("openai_key"):
        st.error("请先在左侧网页 API 设置中填写 OpenAI API Key，视觉识读需要它。")
        st.stop()
    if writer_provider == "未配置":
        st.error("服务器尚未配置教案生成模型API。")
        st.stop()
    if not ppt_file or not docx_file:
        st.error("请同时上传PPT和原教学设计。")
        st.stop()

    try:
        validate_size(ppt_file, MAX_PPT_MB, "PPT")
        validate_size(docx_file, MAX_DOCX_MB, "教学设计")

        ppt_bytes = ppt_file.getvalue()
        docx_bytes = docx_file.getvalue()

        progress = st.progress(0, text="正在读取文件……")

        slides = extract_ppt_text(ppt_bytes)
        template_text, headings = extract_template(docx_bytes)
        if len(headings) < 3:
            raise ValueError(
                "未能可靠识别原教学设计框架。请确认Word中包含“教材分析、教学目标、教学过程”等栏目标题。"
            )

        ppt_text = "\n\n".join(
            f"【PPT第{s['page']}页】\n{s['text'] or '（本页未提取到文本，将依赖视觉识读）'}"
            for s in slides
        )

        progress.progress(20, text="正在云端将PPT转换为PDF……")
        pdf_bytes = pptx_to_pdf(ppt_bytes)

        progress.progress(40, text="正在识读PPT图片、地图、图表和史料……")
        vision_report = analyze_pdf_visual(pdf_bytes, vision_model, st.session_state.openai_key)

        progress.progress(65, text="正在对照原教学设计，删除不匹配内容……")
        writer_prompt = build_writer_prompt(
            template_text, headings, ppt_text, vision_report,
            subject, grade, textbook, periods, extra
        )

        if writer_provider == "DeepSeek":
            raw = call_deepseek(writer_model, writer_prompt, st.session_state.deepseek_key, thinking)
        else:
            raw = call_openai_writer(writer_model, writer_prompt, st.session_state.openai_key)

        progress.progress(85, text="正在重建Word并执行一致性检查……")
        result = safe_json_load(raw)
        final_docx = make_final_docx(
            docx_bytes, result, headings, subject, grade, periods
        )

        progress.progress(100, text="完成")
        st.success("最终教学设计已生成。")

        tabs = st.tabs(["📄 最终教学设计", "🔍 修改核对", "🖼️ PPT视觉报告"])
        section_map = {
            str(x.get("heading", "")).strip(): str(x.get("content", "")).strip()
            for x in result.get("sections", [])
        }

        with tabs[0]:
            st.header(result.get("title") or "教学设计")
            for h in headings:
                st.subheader(h)
                st.markdown(section_map.get(h) or "需教师补充")

        with tabs[1]:
            st.markdown(check_markdown(result))

        with tabs[2]:
            st.markdown(vision_report)

        filename = f"{result.get('title') or '教学设计'}_PPT适配版.docx"
        filename = re.sub(r'[\\/:*?"<>|]', "_", filename)

        st.download_button(
            "⬇️ 下载最终Word教学设计",
            data=final_docx,
            file_name=filename,
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            use_container_width=True,
        )

        st.caption("上传文件仅用于本次生成流程；本程序未配置数据库或永久文件存储。")

    except Exception as e:
        st.error("生成失败")
        st.exception(e)
